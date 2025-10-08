from typing import Dict
from pathlib import Path
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import docx

def parse_pdf(path: str) -> Dict:
    try:
        reader = PdfReader(path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return {'text': text}
    except Exception as e:
        return {'text': f"Error reading PDF: {e}"}

def parse_pptx(path: str) -> Dict:
    try:
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    texts.append(shape.text_frame.text)
                elif hasattr(shape, "text"):
                    texts.append(shape.text)
            slides.append("\n".join(texts))
        return {'text': "\n".join(slides)}
    except Exception as e:
        return {'text': f"Error reading PPTX: {e}"}

def parse_docx(path: str) -> Dict:
    try:
        doc = docx.Document(path)
        paras = [p.text for p in doc.paragraphs]
        return {'text': "\n".join(paras)}
    except Exception as e:
        return {'text': f"Error reading DOCX: {e}"}

def parse_csv(path: str) -> Dict:
    try:
        df = pd.read_csv(path, dtype=str, keep_default_na=False)
        rows = []
        for _, row in df.iterrows():
            rows.append(" | ".join([f"{c}: {row[c]}" for c in df.columns]))
        return {'text': "\n".join(rows)}
    except Exception as e:
        return {'text': f"Error reading CSV: {e}"}

def parse_txt(path: str) -> Dict:
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return {'text': f.read()}
    except Exception as e:
        return {'text': f"Error reading TXT: {e}"}

def parse_file(path: str) -> Dict:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == '.pdf':
        return parse_pdf(path)
    if suffix in ['.pptx', '.ppt']:
        return parse_pptx(path)
    if suffix in ['.docx', '.doc']:
        return parse_docx(path)
    if suffix == '.csv':
        return parse_csv(path)
    if suffix in ['.txt', '.md']:
        return parse_txt(path)
    # Default fallback
    return parse_txt(path)
